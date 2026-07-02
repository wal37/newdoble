from __future__ import annotations

import sys
from pathlib import Path


BASE_DIR = Path(__file__).resolve().parent
for candidate in [
    Path("/Users/wale/Desktop/softtttt/aiapis/apiweb/.venv/lib/python3.13/site-packages"),
    Path("/Users/wale/Desktop/softtttt/foodup/foodweb/.venv/lib/python3.13/site-packages"),
]:
    if candidate.exists() and str(candidate) not in sys.path:
        sys.path.append(str(candidate))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = BASE_DIR / "DOBLEU_PITCH_DECK.pptx"


def rgb(hex_code: str) -> RGBColor:
    hex_code = hex_code.replace("#", "")
    return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))


PAPER = rgb("#F7F1E8")
WARM = rgb("#EEDFCF")
INK = rgb("#111422")
SOFT = rgb("#59617A")
NIGHT = rgb("#161A2F")
NAVY = rgb("#202B57")
VIOLET = rgb("#6B57FF")
ROSE = rgb("#FF8FB1")
CYAN = rgb("#9DE4FF")
GOLD = rgb("#FFCB69")
WHITE = rgb("#FFFFFF")


SLIDES = [
    {
        "type": "cover",
        "title": "Dobleu",
        "subtitle": "Research participation infrastructure for AI labs, user research teams, expert panels, and contributor operations.",
        "eyebrow": "Research Participation Network",
        "stats": [("482", "qualified applicants in pipeline"), ("multi-format", "studies and evaluator tasks"), ("repeatable", "screening and payout ops")],
    },
    {
        "type": "three_panel",
        "title": "Problem",
        "subtitle": "Participant operations are still stitched together with forms, spreadsheets, manual emails, and ad hoc screening.",
        "items": [
            ("Acquisition is fragmented", "Interest forms, sourcing lists, and referrals rarely connect to a reusable participant system."),
            ("Quality is hard to trust", "Teams need faster ways to screen fit, verify expertise, and flag low-quality contributors."),
            ("Study execution stays manual", "Scheduling, notes, payout readiness, and repeat participant management remain operational bottlenecks."),
        ],
    },
    {
        "type": "signal",
        "title": "Solution",
        "subtitle": "Dobleu is a product-led operations layer for participant registration, qualification, study management, and contributor trust.",
        "bullets": [
            "Unified participant intake and registration flows",
            "Screening rules for expertise, demographics, language, and device fit",
            "Study assignment for interviews, evaluations, surveys, and task-based work",
            "Quality history, repeat cohorts, and payout tracking in one system",
        ],
    },
    {
        "type": "three_panel",
        "title": "Why now",
        "subtitle": "AI evaluation demand, human-in-the-loop workflows, and recurring research programs are increasing the cost of weak participant operations.",
        "items": [
            ("More programs need people", "AI labs, product teams, and specialist operators increasingly rely on structured participant feedback and evaluator supply."),
            ("Manual ops do not scale", "Spreadsheet-driven qualification, outreach, and payout workflows create quality risk and operational drag."),
            ("Reuse matters more", "The value of a participant network compounds when teams can trust history, fit, and repeat engagement data."),
        ],
    },
    {
        "type": "map",
        "title": "Platform",
        "subtitle": "One platform connects the full participant research lifecycle.",
        "left_title": "Research teams",
        "left_points": ["AI evaluation programs", "User research and insights", "Expert panels and operations"],
        "center_title": "Dobleu core",
        "center_points": ["Registration", "Screening", "Study ops", "Quality", "Payments"],
        "right_title": "Participants",
        "right_points": ["Applicants", "Qualified contributors", "Repeat evaluators"],
    },
    {
        "type": "stack",
        "title": "Product modules",
        "subtitle": "Built for participant-heavy workflows, not generic CRM usage.",
        "layers": [
            ("Participant profile layer", "Identity fields, expertise, demographics, language, availability, and device attributes."),
            ("Qualification layer", "Screeners, manual review, scoring, cohort rules, and approval logic."),
            ("Study execution layer", "Invites, assignment, scheduling, notes, completion states, and follow-up."),
            ("Trust and payout layer", "Quality flags, contributor history, compensation status, and retention visibility."),
        ],
    },
    {
        "type": "three_panel",
        "title": "Use cases",
        "subtitle": "Dobleu supports different research formats without rebuilding the participant stack every time.",
        "items": [
            ("AI evaluation cohorts", "Recruit multilingual evaluators, experts, and repeat raters for model testing and judgment tasks."),
            ("User research operations", "Run participant pipelines for interviews, product tests, concept validation, and feedback loops."),
            ("Task-based contributor work", "Coordinate annotation, ranking, review, and workflow-specific completion programs."),
        ],
    },
    {
        "type": "timeline",
        "title": "Workflow",
        "subtitle": "Operational flow from interest to qualified contribution and repeat engagement.",
        "steps": [
            ("Register", "Collect participant profile data and structured eligibility details."),
            ("Screen", "Apply fit rules, review logic, and quality gates before assignment."),
            ("Run study", "Invite, schedule, assign, and track completion from one ops layer."),
            ("Retain", "Re-engage high-quality contributors with history and payout accountability."),
        ],
    },
    {
        "type": "ledger",
        "title": "Data system",
        "subtitle": "Each study keeps its operational memory attached to the participant network.",
        "rows": [
            ("Model evaluation", "Bilingual reviewers", "active"),
            ("User interview sprint", "Fintech cohort", "scheduling"),
            ("Expert annotation", "Healthcare specialists", "review"),
            ("UX concept feedback", "Repeat contributors", "paid"),
        ],
    },
    {
        "type": "three_panel",
        "title": "Target market",
        "subtitle": "The initial market centers on organizations that run repeated human-in-the-loop research workflows.",
        "items": [
            ("AI labs and model teams", "Need evaluator supply, structured cohorts, and consistent participant quality."),
            ("User research teams", "Need repeatable participant pipelines instead of one-off recruiting."),
            ("Specialist programs", "Need domain experts, verified contributors, and operational traceability."),
        ],
    },
    {
        "type": "signal",
        "title": "Differentiation",
        "subtitle": "Dobleu is not just lead capture or survey software. It is operational infrastructure for participant quality and reuse.",
        "bullets": [
            "Research-specific intake and screening instead of generic forms",
            "Participant history and trust signals across multiple studies",
            "Task routing, completion, and payout visibility in the same system",
            "Built for repeat cohort operations, not one-off recruitment only",
        ],
    },
    {
        "type": "timeline",
        "title": "Roadmap",
        "subtitle": "The product expands from participant ops into a broader contributor intelligence layer.",
        "steps": [
            ("Phase 1", "Registration, screener flows, and study intake."),
            ("Phase 2", "Scheduling, study dashboards, and repeat cohort logic."),
            ("Phase 3", "Quality scoring, contributor reputation, and payout automation."),
            ("Phase 4", "Deeper AI-assisted matching and enterprise research controls."),
        ],
    },
    {
        "type": "three_panel",
        "title": "Go-to-market",
        "subtitle": "Initial distribution focuses on organizations that already run repeated participant workflows and feel operational pain today.",
        "items": [
            ("Direct enterprise outreach", "Target AI labs, research operations teams, and program owners with recurring evaluator or participant demand."),
            ("Workflow-led adoption", "Land with one study pipeline, then expand into screening, retention, and payout visibility across teams."),
            ("Partner channels", "Support specialist networks, expert programs, and research operators that need branded participant infrastructure."),
        ],
    },
    {
        "type": "signal",
        "title": "Business model",
        "subtitle": "Commercialization can span platform access, cohort operations, and higher-control enterprise workflows.",
        "bullets": [
            "Subscription plans for research teams",
            "Usage-based participant operations or study volumes",
            "Enterprise and API partnerships for large programs",
            "Specialized workflows for expert panels and evaluator networks",
        ],
    },
    {
        "type": "close",
        "title": "Dobleu",
        "subtitle": "A sharper platform for recruiting, qualifying, and operating participant networks at scale.",
        "contact": "contact@dobleu.net",
    },
]


def add_textbox(slide, left, top, width, height, text, size=18, color=INK, bold=False, font="Instrument Sans", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, fill, line=None, radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_cover(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(6.4), Inches(0.2), Inches(5.0), Inches(5.0))
    orb.fill.solid()
    orb.fill.fore_color.rgb = CYAN
    orb.fill.transparency = 0.68
    orb.line.fill.background()
    block = add_card(slide, Inches(0.45), Inches(0.45), Inches(12.4), Inches(6.2), WHITE, line=WARM)
    block.fill.transparency = 0.1
    add_textbox(slide, Inches(0.85), Inches(0.78), Inches(2.6), Inches(0.4), data["eyebrow"], size=16, color=VIOLET, bold=True, font="Space Grotesk")
    add_textbox(slide, Inches(0.85), Inches(1.25), Inches(6.0), Inches(1.4), data["title"], size=28, color=INK, bold=True, font="Space Grotesk")
    add_textbox(slide, Inches(0.85), Inches(2.05), Inches(5.5), Inches(1.2), data["subtitle"], size=17, color=SOFT)
    mosaic_specs = [
        (Inches(7.1), Inches(1.0), Inches(2.2), Inches(1.5), ROSE),
        (Inches(9.45), Inches(1.0), Inches(2.2), Inches(1.5), CYAN),
        (Inches(7.1), Inches(2.7), Inches(2.2), Inches(1.5), GOLD),
        (Inches(9.45), Inches(2.7), Inches(2.2), Inches(1.5), rgb("#D9D4FF")),
    ]
    labels = ["Screeners", "Expert panels", "Study ops", "Payouts"]
    for spec, label in zip(mosaic_specs, labels):
        add_card(slide, *spec, line=WHITE)
        add_textbox(slide, spec[0] + Inches(0.18), spec[1] + Inches(0.36), spec[2] - Inches(0.36), Inches(0.5), label, size=16, color=INK, bold=True, font="Space Grotesk")
    base_left = Inches(0.85)
    for idx, (num, label) in enumerate(data["stats"]):
        x = base_left + Inches(idx * 2.2)
        add_card(slide, x, Inches(5.35), Inches(2.0), Inches(0.9), WHITE, line=WARM)
        add_textbox(slide, x + Inches(0.16), Inches(5.52), Inches(0.9), Inches(0.32), num, size=18, color=NAVY, bold=True, font="Space Grotesk")
        add_textbox(slide, x + Inches(0.88), Inches(5.5), Inches(1.0), Inches(0.35), label, size=9, color=SOFT)


def add_title(slide, title, subtitle, light=False):
    color = WHITE if light else INK
    sub = rgb("#D7DDF3") if light else SOFT
    add_textbox(slide, Inches(0.6), Inches(0.55), Inches(5.7), Inches(0.6), title, size=24, color=color, bold=True, font="Space Grotesk")
    add_textbox(slide, Inches(0.6), Inches(1.05), Inches(8.0), Inches(0.65), subtitle, size=13, color=sub)


def add_three_panel(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    colors = [WHITE, rgb("#F3EEFF"), rgb("#FFF0F4")]
    for idx, (heading, copy) in enumerate(data["items"]):
        left = Inches(0.6 + idx * 4.05)
        add_card(slide, left, Inches(2.0), Inches(3.65), Inches(3.7 if idx == 1 else 3.5), colors[idx], line=WARM)
        add_textbox(slide, left + Inches(0.24), Inches(2.3), Inches(2.9), Inches(0.45), heading, size=17, color=INK, bold=True, font="Space Grotesk")
        add_textbox(slide, left + Inches(0.24), Inches(2.9), Inches(2.95), Inches(1.6), copy, size=12, color=SOFT)


def add_signal(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_card(slide, Inches(0.4), Inches(0.45), Inches(12.5), Inches(6.2), WHITE, line=WARM)
    add_title(slide, data["title"], data["subtitle"])
    for idx, bullet in enumerate(data["bullets"]):
        y = Inches(2.0 + idx * 1.0)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.8), y + Inches(0.08), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = [VIOLET, ROSE, CYAN, GOLD][idx % 4]
        dot.line.fill.background()
        add_textbox(slide, Inches(1.1), y, Inches(6.7), Inches(0.45), bullet, size=16, color=INK, font="Space Grotesk")
    side = add_card(slide, Inches(8.1), Inches(1.7), Inches(4.2), Inches(3.8), NIGHT, line=NIGHT)
    side.fill.transparency = 0.02
    side_labels = ["Intake", "Screen", "Assign", "Retain"]
    for idx, label in enumerate(side_labels):
        y = Inches(2.05 + idx * 0.82)
        add_card(slide, Inches(8.45), y, Inches(3.45), Inches(0.55), [rgb("#2A2E53"), rgb("#4F2F7D"), rgb("#2C4E78"), rgb("#755628")][idx], line=[rgb("#2A2E53"), rgb("#4F2F7D"), rgb("#2C4E78"), rgb("#755628")][idx])
        add_textbox(slide, Inches(8.7), y + Inches(0.13), Inches(2.9), Inches(0.25), label, size=13, color=WHITE, bold=True, font="Space Grotesk")


def add_map(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    cols = [(Inches(0.75), data["left_title"], data["left_points"], rgb("#FFF0F4")),
            (Inches(4.45), data["center_title"], data["center_points"], rgb("#EEF2FF")),
            (Inches(8.15), data["right_title"], data["right_points"], rgb("#ECFBFF"))]
    for left, title, points, fill in cols:
        add_card(slide, left, Inches(2.0), Inches(3.15), Inches(3.7), fill, line=WARM)
        add_textbox(slide, left + Inches(0.24), Inches(2.28), Inches(2.8), Inches(0.4), title, size=17, color=INK, bold=True, font="Space Grotesk")
        for idx, point in enumerate(points):
            add_textbox(slide, left + Inches(0.28), Inches(2.85 + idx * 0.58), Inches(2.85), Inches(0.35), f"• {point}", size=12, color=SOFT)


def add_stack(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, data["title"], data["subtitle"])
    fills = [rgb("#F0EDFF"), rgb("#FCEAF0"), rgb("#E9F8FF"), rgb("#FFF3D8")]
    for idx, (heading, copy) in enumerate(data["layers"]):
        y = Inches(1.9 + idx * 1.05)
        add_card(slide, Inches(0.7), y, Inches(11.5), Inches(0.82), fills[idx], line=WARM)
        add_textbox(slide, Inches(0.95), y + Inches(0.16), Inches(2.7), Inches(0.28), heading, size=14, color=INK, bold=True, font="Space Grotesk")
        add_textbox(slide, Inches(3.35), y + Inches(0.15), Inches(8.2), Inches(0.34), copy, size=11, color=SOFT)


def add_timeline(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NIGHT
    add_title(slide, data["title"], data["subtitle"], light=True)
    for idx, (heading, copy) in enumerate(data["steps"]):
        left = Inches(0.6 + idx * 3.05)
        add_card(slide, left, Inches(2.05), Inches(2.75), Inches(3.2), rgb("#262B4B"), line=rgb("#2F355B"))
        add_textbox(slide, left + Inches(0.2), Inches(2.3), Inches(0.45), Inches(0.28), str(idx + 1), size=14, color=WHITE, bold=True, font="Space Grotesk")
        add_textbox(slide, left + Inches(0.62), Inches(2.28), Inches(1.9), Inches(0.35), heading, size=16, color=WHITE, bold=True, font="Space Grotesk")
        add_textbox(slide, left + Inches(0.22), Inches(2.95), Inches(2.25), Inches(1.25), copy, size=11, color=rgb("#D3DCF9"))


def add_ledger(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    add_card(slide, Inches(0.65), Inches(1.9), Inches(11.7), Inches(3.95), WHITE, line=WARM)
    add_textbox(slide, Inches(0.95), Inches(2.15), Inches(2.5), Inches(0.25), "Study type", size=11, color=SOFT, bold=True)
    add_textbox(slide, Inches(5.3), Inches(2.15), Inches(2.0), Inches(0.25), "Cohort", size=11, color=SOFT, bold=True)
    add_textbox(slide, Inches(9.35), Inches(2.15), Inches(1.4), Inches(0.25), "Status", size=11, color=SOFT, bold=True)
    for idx, row in enumerate(data["rows"]):
        y = Inches(2.6 + idx * 0.75)
        add_textbox(slide, Inches(0.95), y, Inches(3.8), Inches(0.3), row[0], size=14, color=INK, bold=True, font="Space Grotesk")
        add_textbox(slide, Inches(5.3), y, Inches(2.8), Inches(0.3), row[1], size=12, color=SOFT)
        pill = add_card(slide, Inches(9.25), y - Inches(0.03), Inches(1.55), Inches(0.4), rgb("#EEE8FF"), line=rgb("#EEE8FF"))
        pill.fill.transparency = 0.02
        add_textbox(slide, Inches(9.47), y + Inches(0.05), Inches(1.0), Inches(0.2), row[2], size=10, color=VIOLET, bold=True)


def add_close(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.0), Inches(0.4), Inches(4.0), Inches(4.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = VIOLET
    glow.fill.transparency = 0.62
    glow.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.65), data["title"], size=28, color=WHITE, bold=True, font="Space Grotesk")
    add_textbox(slide, Inches(0.8), Inches(2.3), Inches(6.0), Inches(1.2), data["subtitle"], size=17, color=rgb("#D5D9F4"))
    add_card(slide, Inches(0.8), Inches(4.85), Inches(3.5), Inches(0.75), rgb("#262B4B"), line=rgb("#262B4B"))
    add_textbox(slide, Inches(1.05), Inches(5.05), Inches(2.7), Inches(0.25), data["contact"], size=14, color=WHITE, bold=True, font="Space Grotesk")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Dobleu Pitch Deck"

    layout = prs.slide_layouts[6]
    for entry in SLIDES:
        slide = prs.slides.add_slide(layout)
        if entry["type"] == "cover":
            add_cover(slide, entry)
        elif entry["type"] == "three_panel":
            add_three_panel(slide, entry)
        elif entry["type"] == "signal":
            add_signal(slide, entry)
        elif entry["type"] == "map":
            add_map(slide, entry)
        elif entry["type"] == "stack":
            add_stack(slide, entry)
        elif entry["type"] == "timeline":
            add_timeline(slide, entry)
        elif entry["type"] == "ledger":
            add_ledger(slide, entry)
        elif entry["type"] == "close":
            add_close(slide, entry)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
